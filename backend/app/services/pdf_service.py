import fitz  # PyMuPDF
import base64
import uuid
import json
import os
import threading
from typing import Dict, Optional, List
from collections import OrderedDict
from app.models.pdf import PdfInfo, PageRender, ThumbnailRender, PdfOutlineItem, PageSize, Annotation
from app.core.config import settings


class PasswordRequiredError(Exception):
    """fitz no define PasswordError; esta señala al router que devuelva 401."""


class PdfService:
    def __init__(self):
        self._docs: Dict[str, fitz.Document] = {}
        self._infos: Dict[str, PdfInfo] = {}
        self._dirty: Dict[str, bool] = {}  # Track unsaved changes
        self._passwords: Dict[str, Optional[str]] = {}  # Kept so evicted protected docs can be reopened
        self._render_cache: OrderedDict = OrderedDict()  # (doc_id, page_num, zoom) -> PageRender
        self._thumb_cache: OrderedDict = OrderedDict()   # (doc_id, page_num) -> ThumbnailRender
        self._render_cache_max = 150
        self._thumb_cache_max = 300
        self._snap_cache: OrderedDict = OrderedDict()  # (doc_id, page_num) -> List[dict]
        self._snap_cache_max = 60
        # LRU eviction: with 60+ plans open the backend would otherwise keep every
        # fitz.Document (and its cached page objects) alive. We close the least-recently
        # used *non-dirty* documents and transparently reopen them from disk on access.
        # Dirty docs are never evicted (reopening would lose unsaved edits).
        self._lru: "OrderedDict[str, None]" = OrderedDict()  # doc_id -> None, oldest first
        self._max_live_docs = 12
        # PyMuPDF is not thread-safe and handlers now run in FastAPI's threadpool,
        # so serialize all access to documents and caches.
        self._lock = threading.RLock()

    def _acquire(self, doc_id: str) -> Optional[fitz.Document]:
        """Return the live fitz.Document for doc_id, reopening it from disk if it was
        evicted. Updates LRU recency and evicts other inactive non-dirty docs.
        Returns None for unknown docs or if the file can't be reopened."""
        with self._lock:
            if doc_id not in self._infos:
                return None
            doc = self._docs.get(doc_id)
            if doc is None:
                info = self._infos[doc_id]
                try:
                    doc = fitz.open(info.file_path)
                except Exception:
                    return None
                pw = self._passwords.get(doc_id)
                if doc.needs_pass and pw:
                    doc.authenticate(pw)
                self._docs[doc_id] = doc
                self._dirty.setdefault(doc_id, False)
            self._lru.pop(doc_id, None)
            self._lru[doc_id] = None  # most-recently used
            self._evict_inactive()
            return self._docs.get(doc_id)

    def _evict_inactive(self):
        """Close least-recently-used non-dirty documents when over the live cap.
        Must be called holding self._lock."""
        if len(self._docs) <= self._max_live_docs:
            return
        for did in list(self._lru.keys()):
            if len(self._docs) <= self._max_live_docs:
                break
            if did not in self._docs:
                self._lru.pop(did, None)
                continue
            if self._dirty.get(did):
                continue  # keep unsaved work in memory
            doc = self._docs.pop(did, None)
            if doc is not None:
                try:
                    doc.close()
                except Exception:
                    pass
            self._lru.pop(did, None)

    @staticmethod
    def _capped_scale(page, desired_scale: float, max_px: int) -> float:
        """Clamp the render scale so very large pages (e.g. architectural drawings)
        don't produce multi-thousand-pixel bitmaps that block the render thread."""
        md = max(page.rect.width, page.rect.height)
        if md <= 0:
            return desired_scale
        return max(0.05, min(desired_scale, max_px / md))

    def _invalidate_render_cache(self, doc_id: str):
        keys_to_remove = [k for k in self._render_cache.keys() if k[0] == doc_id]
        for k in keys_to_remove:
            del self._render_cache[k]
        keys_to_remove = [k for k in self._thumb_cache.keys() if k[0] == doc_id]
        for k in keys_to_remove:
            del self._thumb_cache[k]
        keys_to_remove = [k for k in self._snap_cache.keys() if k[0] == doc_id]
        for k in keys_to_remove:
            del self._snap_cache[k]

    def open_document(self, file_path: str, password: Optional[str] = None) -> PdfInfo:
        doc = fitz.open(file_path)
        if doc.needs_pass:
            if not password:
                doc.close()
                raise PasswordRequiredError("Document requires a password")
            auth_result = doc.authenticate(password)
            if not auth_result:
                doc.close()
                raise PasswordRequiredError("Incorrect password")
        doc_id = str(uuid.uuid4())

        page_sizes = []
        for i in range(len(doc)):
            rect = doc.load_page(i).rect
            page_sizes.append(PageSize(page_num=i, width=rect.width, height=rect.height))

        info = PdfInfo(
            doc_id=doc_id,
            file_path=file_path,
            page_count=len(doc),
            title=doc.metadata.get("title"),
            author=doc.metadata.get("author"),
            subject=doc.metadata.get("subject"),
            current_page=0,
            page_sizes=page_sizes
        )

        with self._lock:
            self._docs[doc_id] = doc
            self._infos[doc_id] = info
            self._dirty[doc_id] = False
            self._passwords[doc_id] = password
            self._lru[doc_id] = None
            self._evict_inactive()
        return info

    def get_document(self, doc_id: str) -> Optional[fitz.Document]:
        return self._acquire(doc_id)

    def get_info(self, doc_id: str) -> Optional[PdfInfo]:
        return self._infos.get(doc_id)

    def get_page_info(self, doc_id: str, page_num: int, zoom: float = 1.0) -> Optional[dict]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None
            page = doc.load_page(page_num)
            max_px = 3000 if zoom <= 1.6 else 6000
            scale = self._capped_scale(page, zoom * settings.RENDER_DPI / 72, max_px)
            pw = int(page.rect.width * scale)
            ph = int(page.rect.height * scale)
            return {
                "page_num": page_num,
                "width": pw,
                "height": ph,
                "original_width": page.rect.width,
                "original_height": page.rect.height,
            }

    def get_page_image_bytes(self, doc_id: str, page_num: int, zoom: float = 1.0) -> Optional[bytes]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None
            cache_key = (doc_id, page_num, zoom)
            if cache_key in self._render_cache:
                render = self._render_cache[cache_key]
                # Extract raw base64 bytes
                b64 = render.image_base64.split(',')[1]
                return base64.b64decode(b64)
            page = doc.load_page(page_num)
            max_px = 3000 if zoom <= 1.6 else 6000
            scale = self._capped_scale(page, zoom * settings.RENDER_DPI / 72, max_px)
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            return pix.tobytes("png")

    def render_tile_bytes(self, doc_id: str, page_num: int, x0: float, y0: float,
                          x1: float, y1: float, zoom: float) -> Optional[bytes]:
        """Render only a sub-rectangle (in PDF points) of a page at the true target
        zoom. Used for crisp deep-zoom on dense plans without rasterizing the whole
        page at huge resolution. The tile pixel size is capped so a single request
        can't block the render thread."""
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None
            page = doc.load_page(page_num)
            clip = fitz.Rect(x0, y0, x1, y1) & page.rect
            if clip.is_empty or clip.width <= 0 or clip.height <= 0:
                return None
            desired = zoom * settings.RENDER_DPI / 72
            md = max(clip.width, clip.height)
            scale = max(0.05, min(desired, 4000 / md)) if md > 0 else desired
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
            return pix.tobytes("png")

    def render_page(self, doc_id: str, page_num: int, zoom: float = 1.0) -> Optional[PageRender]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None

            cache_key = (doc_id, page_num, zoom)
            if cache_key in self._render_cache:
                return self._render_cache[cache_key]

            page = doc.load_page(page_num)
            # Modest cap for the base view (fast first paint), larger cap when the user
            # zooms in (sharper on demand) — keeps huge CAD plans both responsive and legible.
            max_px = 3000 if zoom <= 1.6 else 6000
            scale = self._capped_scale(page, zoom * settings.RENDER_DPI / 72, max_px)
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img_data = pix.tobytes("png")
            img_b64 = base64.b64encode(img_data).decode('utf-8')

            result = PageRender(
                page_num=page_num,
                image_base64=f"data:image/png;base64,{img_b64}",
                width=pix.width,
                height=pix.height,
                original_width=page.rect.width,
                original_height=page.rect.height
            )
            self._render_cache[cache_key] = result
            if len(self._render_cache) > self._render_cache_max:
                self._render_cache.popitem(last=False)
            return result

    def render_thumbnail(self, doc_id: str, page_num: int) -> Optional[ThumbnailRender]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None

            cache_key = (doc_id, page_num)
            if cache_key in self._thumb_cache:
                return self._thumb_cache[cache_key]

            page = doc.load_page(page_num)
            scale = self._capped_scale(page, settings.THUMBNAIL_DPI / 72, 300)
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat, alpha=False)

            img_data = pix.tobytes("png")
            img_b64 = base64.b64encode(img_data).decode('utf-8')

            result = ThumbnailRender(
                page_num=page_num,
                image_base64=f"data:image/png;base64,{img_b64}",
                width=pix.width,
                height=pix.height
            )
            self._thumb_cache[cache_key] = result
            if len(self._thumb_cache) > self._thumb_cache_max:
                self._thumb_cache.popitem(last=False)
            return result

    def get_outline(self, doc_id: str) -> List[PdfOutlineItem]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return []
            outline = doc.get_toc()
        result: List[PdfOutlineItem] = []
        stack: List[tuple] = []
        
        for item in outline:
            level, title, page = item[0], item[1], item[2] - 1
            node = PdfOutlineItem(title=title, page=page, children=[])
            
            while len(stack) >= level:
                stack.pop()
            
            if not stack:
                result.append(node)
            else:
                parent = stack[-1][1]
                if parent.children is None:
                    parent.children = []
                parent.children.append(node)
            
            stack.append((level, node))
        
        return result

    def search_text(self, doc_id: str, query: str, limit: int = 500) -> List[dict]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return []

            results = []
            for page_num in range(len(doc)):
                if len(results) >= limit:
                    break
                page = doc.load_page(page_num)
                rects = page.search_for(query)
                for rect in rects:
                    if len(results) >= limit:
                        break
                    band = fitz.Rect(0, rect.y0 - 1, page.rect.width, rect.y1 + 1)
                    snippet = page.get_textbox(band).replace("\n", " ").strip()
                    if len(snippet) > 160:
                        snippet = snippet[:160] + "…"
                    results.append({
                        "page": page_num,
                        "x": rect.x0,
                        "y": rect.y0,
                        "width": rect.width,
                        "height": rect.height,
                        "snippet": snippet,
                    })
            return results

    # --- Document operations ---

    def rotate_page(self, doc_id: str, page_num: int, degrees: int) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return False
        page = doc.load_page(page_num)
        page.set_rotation((page.rotation + degrees) % 360)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def rotate_all_pages(self, doc_id: str, degrees: int) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        for i in range(len(doc)):
            page = doc.load_page(i)
            page.set_rotation((page.rotation + degrees) % 360)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def rotate_pages(self, doc_id: str, pages: List[int], degrees: int) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        for p in pages:
            if 0 <= p < len(doc):
                page = doc.load_page(p)
                page.set_rotation((page.rotation + degrees) % 360)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def reorder_pages(self, doc_id: str, new_order: List[int]) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        if len(new_order) != len(doc) or set(new_order) != set(range(len(doc))):
            return False
        doc.select(new_order)
        info = self._infos.get(doc_id)
        if info:
            info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def add_watermark(self, doc_id: str, text: str, color: str = "#888888", fontsize: float = 48, angle: int = 45, opacity: float = 0.3) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255.0 for i in (0, 2, 4)) if color.startswith('#') else (0.5, 0.5, 0.5)
        for i in range(len(doc)):
            page = doc.load_page(i)
            rect = page.rect
            # Approximate centering for rotated text
            text_width = fitz.get_text_length(text, fontsize=fontsize)
            insert_x = rect.width / 2 - text_width / 4
            insert_y = rect.height / 2 + fontsize / 4
            page.insert_text((insert_x, insert_y), text, fontsize=fontsize, color=rgb, rotate=angle, overlay=True)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def create_blank_pdf(self, output_path: str, page_width: float = 595, page_height: float = 842, page_count: int = 1) -> Optional[PdfInfo]:
        doc = fitz.open()
        for _ in range(page_count):
            doc.new_page(width=page_width, height=page_height)
        doc.save(output_path)
        doc.close()
        return self.open_document(output_path)

    def redact_area(self, doc_id: str, page_num: int, x: float, y: float, width: float, height: float) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return False
        page = doc.load_page(page_num)
        rect = fitz.Rect(x, y, x + width, y + height)
        page.add_redact_annot(rect)
        page.apply_redactions()
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def crop_page(self, doc_id: str, page_num: int, top: float, right: float, bottom: float, left: float) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return False
        page = doc.load_page(page_num)
        rect = page.rect
        new_rect = fitz.Rect(rect.x0 + left, rect.y0 + top, rect.x1 - right, rect.y1 - bottom)
        if new_rect.width <= 0 or new_rect.height <= 0:
            return False
        page.set_cropbox(new_rect)
        info = self._infos.get(doc_id)
        if info:
            info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def export_excel(self, doc_id: str, output_path: str) -> bool:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return False
            try:
                from openpyxl import Workbook
                wb = Workbook()
                first = True

                def next_sheet(title: str):
                    nonlocal first
                    ws = wb.active if first else wb.create_sheet()
                    ws.title = title[:31]
                    first = False
                    return ws

                for i in range(len(doc)):
                    page = doc.load_page(i)
                    tables = []
                    try:
                        tables = page.find_tables().tables
                    except Exception:
                        tables = []
                    if tables:
                        # Preserve real table structure, one sheet per table
                        for t_idx, tab in enumerate(tables):
                            ws = next_sheet(f"P{i+1}_T{t_idx+1}")
                            for row in tab.extract():
                                ws.append([("" if c is None else c) for c in row])
                    else:
                        # Fall back to raw text lines for pages without tables
                        ws = next_sheet(f"Página {i+1}")
                        for line in page.get_text().splitlines()[:1000]:
                            ws.append([line])
                wb.save(output_path)
                return True
            except Exception:
                return False

    def ocr_available(self) -> bool:
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            return True
        except Exception:
            return False

    def make_searchable(self, doc_id: str, pages: Optional[List[int]] = None) -> int:
        """OCR scanned pages and embed an invisible text layer so they become
        searchable/selectable. Returns words added, or -1 if Tesseract is unavailable."""
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return 0
            try:
                import pytesseract
                from pytesseract import Output
                from PIL import Image
                from io import BytesIO
            except Exception:
                return -1
            try:
                pytesseract.get_tesseract_version()
            except Exception:
                return -1

            DPI = 200
            scale = 72.0 / DPI
            indices = pages if pages is not None else list(range(len(doc)))
            total = 0
            for i in indices:
                if i < 0 or i >= len(doc):
                    continue
                page = doc.load_page(i)
                if page.get_text().strip():
                    continue  # already has a text layer
                pix = page.get_pixmap(dpi=DPI)
                img = Image.open(BytesIO(pix.tobytes("png")))
                data = pytesseract.image_to_data(img, lang="spa+eng", output_type=Output.DICT)
                for j, txt in enumerate(data["text"]):
                    if not txt or not txt.strip():
                        continue
                    x = data["left"][j] * scale
                    y = data["top"][j] * scale
                    h = max(4, data["height"][j] * scale)
                    try:
                        page.insert_text((x, y + h), txt, fontsize=h * 0.9, render_mode=3)
                        total += 1
                    except Exception:
                        continue
            if total:
                self._dirty[doc_id] = True
                self._invalidate_render_cache(doc_id)
            return total

    def export_pptx(self, doc_id: str, output_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from io import BytesIO
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=150)
                img_stream = BytesIO(pix.tobytes("png"))
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            prs.save(output_path)
            return True
        except Exception:
            return False

    def ocr_page(self, doc_id: str, page_num: int) -> Optional[str]:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return None
        try:
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=300)
            from io import BytesIO
            from PIL import Image
            import pytesseract
            img = Image.open(BytesIO(pix.tobytes("png")))
            return pytesseract.image_to_string(img, lang='spa+eng')
        except Exception:
            return None

    def get_page_text_data(self, doc_id: str, page_num: int) -> Optional[dict]:
        """Both the positioned text blocks (for the selectable overlay) and the plain
        text (for copy / read-aloud) in a single locked call."""
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None
            page = doc.load_page(page_num)
            blocks = []
            try:
                for b in page.get_text("blocks"):
                    x0, y0, x1, y1, text, *_ = b
                    blocks.append({"x": x0, "y": y0, "width": x1 - x0, "height": y1 - y0, "text": text.strip()})
            except Exception:
                blocks = []
            return {"blocks": blocks, "text": page.get_text()}

    def delete_pages(self, doc_id: str, pages: List[int]) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        # Delete in reverse order to keep indices valid
        for p in sorted(pages, reverse=True):
            if 0 <= p < len(doc):
                doc.delete_page(p)
        info = self._infos.get(doc_id)
        if info:
            info.page_count = len(doc)
            info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def merge_pdf(self, doc_id: str, source_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc or not os.path.exists(source_path):
            return False
        try:
            src = fitz.open(source_path)
            doc.insert_pdf(src)
            src.close()
            info = self._infos.get(doc_id)
            if info:
                info.page_count = len(doc)
                info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
            self._dirty[doc_id] = True
            self._invalidate_render_cache(doc_id)
            return True
        except Exception:
            return False

    def split_pages(self, doc_id: str, pages: List[int], output_path: Optional[str] = None) -> Optional[str]:
        doc = self._acquire(doc_id)
        if not doc:
            return None
        new_doc = fitz.open()
        for p in pages:
            if 0 <= p < len(doc):
                new_doc.insert_pdf(doc, from_page=p, to_page=p)
        save_path = output_path or os.path.join(os.path.dirname(doc.name), f"split_{uuid.uuid4().hex[:8]}.pdf")
        new_doc.save(save_path)
        new_doc.close()
        return save_path

    def compress(self, doc_id: str, output_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        try:
            doc.save(output_path, garbage=4, deflate=True, clean=True)
            return True
        except Exception:
            return False

    def save(self, doc_id: str, output_path: Optional[str] = None) -> Optional[str]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return None
            save_path = output_path or doc.name
            in_place = os.path.abspath(save_path) == os.path.abspath(doc.name)
            temp_path = None
            try:
                import tempfile
                dir_name = os.path.dirname(os.path.abspath(save_path))
                fd, temp_path = tempfile.mkstemp(suffix='.pdf', dir=dir_name)
                os.close(fd)
                doc.save(temp_path, garbage=4, deflate=True)
                if in_place:
                    # On Windows the open fitz.Document holds the file handle, so
                    # os.replace onto doc.name fails with PermissionError. Close,
                    # replace, then reopen from the saved file.
                    doc.close()
                    self._docs.pop(doc_id, None)
                    os.replace(temp_path, save_path)
                    self._dirty[doc_id] = False
                    self._acquire(doc_id)  # reopen so the doc stays usable
                else:
                    os.replace(temp_path, save_path)
                    self._dirty[doc_id] = False
                return save_path
            except Exception:
                try:
                    if temp_path and os.path.exists(temp_path):
                        os.remove(temp_path)
                except Exception:
                    pass
                return None

    def insert_text(self, doc_id: str, page_num: int, x: float, y: float, text: str, color: str = "#000000", fontsize: float = 12) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return False
        page = doc.load_page(page_num)
        # Parse hex color
        rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255.0 for i in (0, 2, 4)) if color.startswith('#') else (0, 0, 0)
        page.insert_text((x, y), text, fontsize=fontsize, color=rgb)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def insert_image(self, doc_id: str, page_num: int, x: float, y: float, width: float, height: float, image_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc) or not os.path.exists(image_path):
            return False
        page = doc.load_page(page_num)
        rect = fitz.Rect(x, y, x + width, y + height)
        page.insert_image(rect, filename=image_path)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def get_snap_points(self, doc_id: str, page_num: int) -> Optional[List[dict]]:
        """Puntos de ajuste para mediciones: extremos y puntos medios de líneas,
        esquinas de rectángulos y extremos de curvas del contenido vectorial.
        Coordenadas en unidades PDF (las mismas que usan las anotaciones)."""
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return None
            cache_key = (doc_id, page_num)
            if cache_key in self._snap_cache:
                self._snap_cache.move_to_end(cache_key)
                return self._snap_cache[cache_key]

            MAX_POINTS = 20000
            points: List[dict] = []
            seen = set()

            def add(x: float, y: float):
                key = (round(x, 1), round(y, 1))
                if key not in seen:
                    seen.add(key)
                    points.append({"x": key[0], "y": key[1]})

            page = doc.load_page(page_num)
            try:
                drawings = page.get_drawings()
            except Exception:
                drawings = []
            for d in drawings:
                if len(points) >= MAX_POINTS:
                    break
                for item in d.get("items", []):
                    kind = item[0]
                    if kind == "l":
                        p1, p2 = item[1], item[2]
                        add(p1.x, p1.y)
                        add(p2.x, p2.y)
                        add((p1.x + p2.x) / 2, (p1.y + p2.y) / 2)
                    elif kind == "re":
                        r = item[1]
                        add(r.x0, r.y0)
                        add(r.x1, r.y0)
                        add(r.x0, r.y1)
                        add(r.x1, r.y1)
                    elif kind == "c":
                        add(item[1].x, item[1].y)
                        add(item[4].x, item[4].y)
                    elif kind == "qu":
                        q = item[1]
                        for p in (q.ul, q.ur, q.ll, q.lr):
                            add(p.x, p.y)
                    if len(points) >= MAX_POINTS:
                        break

            self._snap_cache[cache_key] = points
            while len(self._snap_cache) > self._snap_cache_max:
                self._snap_cache.popitem(last=False)
            return points

    # --- Annotations persistence ---

    def _get_annotations_path(self, file_path: str) -> str:
        return file_path + ".pdfmaster.json"

    def load_annotations(self, doc_id: str) -> List[Annotation]:
        info = self._infos.get(doc_id)
        if not info:
            return []
        path = self._get_annotations_path(info.file_path)
        if not os.path.exists(path):
            return []
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return [Annotation(**ann) for ann in data.get('annotations', [])]
        except Exception:
            return []

    def save_annotations(self, doc_id: str, annotations: List[Annotation]) -> bool:
        info = self._infos.get(doc_id)
        if not info:
            return False
        path = self._get_annotations_path(info.file_path)
        try:
            with open(path, 'w', encoding='utf-8') as f:
                json.dump({"version": 1, "annotations": [ann.model_dump() for ann in annotations]}, f, indent=2)
            return True
        except Exception:
            return False

    def embed_annotations(self, doc_id: str, annotations: List[Annotation]) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        self._invalidate_render_cache(doc_id)
        
        def hex_to_rgb(color: Optional[str]):
            if not color or not color.startswith('#'):
                return (0, 0, 0)
            try:
                return tuple(int(color.lstrip('#')[i:i+2], 16) / 255.0 for i in (0, 2, 4))
            except Exception:
                return (0, 0, 0)
        
        for ann in annotations:
            if ann.page < 0 or ann.page >= len(doc):
                continue
            page = doc.load_page(ann.page)
            color = hex_to_rgb(ann.color)
            
            if ann.type == 'highlight':
                rect = fitz.Rect(ann.x, ann.y, ann.x + (ann.width or 0), ann.y + (ann.height or 0))
                annot = page.add_highlight_annot(rect)
                if annot:
                    annot.set_colors(stroke=color)
                    annot.update()
            elif ann.type == 'underline':
                rect = fitz.Rect(ann.x, ann.y, ann.x + (ann.width or 0), ann.y + (ann.height or 0))
                annot = page.add_underline_annot(rect)
                if annot:
                    annot.set_colors(stroke=color)
                    annot.update()
            elif ann.type == 'strikethrough':
                rect = fitz.Rect(ann.x, ann.y, ann.x + (ann.width or 0), ann.y + (ann.height or 0))
                annot = page.add_strikeout_annot(rect)
                if annot:
                    annot.set_colors(stroke=color)
                    annot.update()
            elif ann.type == 'rect':
                rect = fitz.Rect(ann.x, ann.y, ann.x + (ann.width or 0), ann.y + (ann.height or 0))
                page.draw_rect(rect, color=color, width=ann.lineWidth or 2)
            elif ann.type == 'circle':
                rect = fitz.Rect(ann.x, ann.y, ann.x + (ann.width or 0), ann.y + (ann.height or 0))
                page.draw_oval(rect, color=color, width=ann.lineWidth or 2)
            elif ann.type == 'arrow':
                x1, y1 = ann.x, ann.y
                x2, y2 = ann.x + (ann.width or 0), ann.y + (ann.height or 0)
                page.draw_line(fitz.Point(x1, y1), fitz.Point(x2, y2), color=color, width=ann.lineWidth or 2)
                # Arrowhead
                angle = fitz.utils.degrees(fitz.utils.atan2(y2 - y1, x2 - x1))
                head_len = 10
                p1 = fitz.Point(x2, y2)
                p2 = fitz.Point(x2 - head_len * fitz.utils.cos(angle - 30), y2 - head_len * fitz.utils.sin(angle - 30))
                p3 = fitz.Point(x2 - head_len * fitz.utils.cos(angle + 30), y2 - head_len * fitz.utils.sin(angle + 30))
                page.draw_polygon([p1, p2, p3], color=color, fill=color)
            elif ann.type == 'draw':
                if ann.points and len(ann.points) > 1:
                    pts = [fitz.Point(p.x, p.y) for p in ann.points]
                    page.draw_polyline(pts, color=color, width=ann.lineWidth or 2)
            elif ann.type == 'signature':
                if ann.points and len(ann.points) > 1:
                    pts = [fitz.Point(p.x, p.y) for p in ann.points]
                    page.draw_polyline(pts, color=(0, 0, 0), width=3)
            elif ann.type == 'text':
                fs = ann.fontSize or 14
                # ann.y is the text box top-left; place each line's baseline below it.
                for i, line in enumerate((ann.text or '').split('\n')):
                    page.insert_text((ann.x, ann.y + fs * 0.8 + i * fs * 1.3), line, fontsize=fs, color=color)
            elif ann.type == 'note':
                annot = page.add_text_annot(fitz.Point(ann.x, ann.y), ann.text or 'Nota')
                if annot:
                    annot.set_colors(stroke=color)
                    annot.update()
        
        self._dirty[doc_id] = True
        return True

    def add_header_footer(self, doc_id: str, header: Optional[str] = None, footer: Optional[str] = None, fontsize: float = 10, color: str = "#000000") -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255.0 for i in (0, 2, 4)) if color.startswith('#') else (0, 0, 0)
        for i in range(len(doc)):
            page = doc.load_page(i)
            rect = page.rect
            if header:
                tw = fitz.get_text_length(header, fontsize=fontsize)
                page.insert_text((rect.width / 2 - tw / 2, 20), header, fontsize=fontsize, color=rgb, overlay=True)
            if footer:
                tw = fitz.get_text_length(footer, fontsize=fontsize)
                page.insert_text((rect.width / 2 - tw / 2, rect.height - 10), footer, fontsize=fontsize, color=rgb, overlay=True)
        self._dirty[doc_id] = True
        self._invalidate_render_cache(doc_id)
        return True

    def replace_text(self, doc_id: str, query: str, replace: str, page_num: Optional[int] = None, case_sensitive: bool = False, replace_all: bool = True) -> int:
        doc = self._acquire(doc_id)
        if not doc or not query:
            return 0
        flags = 0 if case_sensitive else fitz.TEXT_DEHYPHENATE
        count = 0
        pages_to_search = [page_num] if page_num is not None else range(len(doc))
        for p in pages_to_search:
            if p < 0 or p >= len(doc):
                continue
            page = doc.load_page(p)
            rects = page.search_for(query, flags=flags)
            if not rects:
                continue
            # Redact old text
            for rect in rects:
                page.add_redact_annot(rect)
            page.apply_redactions()
            # Insert new text at first rect position with default styling
            if rects:
                # Estimate font size from rect height
                est_fontsize = max(8, min(72, int(rects[0].height * 0.8)))
                page.insert_text(rects[0].tl, replace, fontsize=est_fontsize, color=(0, 0, 0), overlay=True)
                count += len(rects)
            if not replace_all:
                break
        if count > 0:
            self._dirty[doc_id] = True
            self._invalidate_render_cache(doc_id)
        return count

    def set_metadata(self, doc_id: str, title: Optional[str] = None, author: Optional[str] = None, subject: Optional[str] = None, keywords: Optional[str] = None) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        meta = doc.metadata
        if title is not None:
            meta['title'] = title
        if author is not None:
            meta['author'] = author
        if subject is not None:
            meta['subject'] = subject
        if keywords is not None:
            meta['keywords'] = keywords
        doc.set_metadata(meta)
        self._dirty[doc_id] = True
        return True

    def get_page_spans(self, doc_id: str, page_num: int) -> List[dict]:
        """Per-span text with bounding boxes, used to build a selectable text layer."""
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return []
            page = doc.load_page(page_num)
            data = page.get_text("dict")
            spans: List[dict] = []
            for block in data.get("blocks", []):
                for line in block.get("lines", []):
                    for sp in line.get("spans", []):
                        text = sp.get("text", "")
                        if not text.strip():
                            continue
                        x0, y0, x1, y1 = sp["bbox"]
                        spans.append({
                            "text": text,
                            "x0": x0, "y0": y0, "x1": x1, "y1": y1,
                            "size": sp.get("size", y1 - y0),
                        })
            return spans

    def export_word(self, doc_id: str) -> Optional[dict]:
        doc = self._acquire(doc_id)
        if not doc:
            return None
        try:
            from docx import Document
            from docx.shared import Pt
            import io
            import base64
            document = Document()
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    for line in text.split('\n'):
                        if line.strip():
                            p = document.add_paragraph(line.strip())
                            p.style.font.size = Pt(11)
                if page_num < len(doc) - 1:
                    document.add_page_break()
            buffer = io.BytesIO()
            document.save(buffer)
            buffer.seek(0)
            data = buffer.read()
            return {
                "filename": os.path.basename(doc.name).replace('.pdf', '.docx'),
                "data_base64": base64.b64encode(data).decode('utf-8'),
            }
        except Exception as e:
            print("Export word error:", e)
            return None

    def get_form_fields(self, doc_id: str, page_num: int) -> List[dict]:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return []
        page = doc.load_page(page_num)
        widgets = page.widgets()
        results = []
        for widget in widgets:
            results.append({
                "field_name": widget.field_name,
                "field_type": widget.field_type_string,
                "field_flags": widget.field_flags,
                "rect": {
                    "x": widget.rect.x0,
                    "y": widget.rect.y0,
                    "width": widget.rect.width,
                    "height": widget.rect.height,
                },
                "value": widget.field_value or "",
                "options": widget.choice_values or [],
            })
        return results

    def set_form_field(self, doc_id: str, page_num: int, field_name: str, value: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return False
        page = doc.load_page(page_num)
        widgets = page.widgets()
        for widget in widgets:
            if widget.field_name == field_name:
                widget.field_value = value
                widget.update()
                self._dirty[doc_id] = True
                return True
        return False

    def get_text_clip(self, doc_id: str, page_num: int, x: float, y: float, w: float, h: float) -> str:
        doc = self._acquire(doc_id)
        if not doc or page_num < 0 or page_num >= len(doc):
            return ""
        page = doc.load_page(page_num)
        rect = fitz.Rect(x, y, x + w, y + h)
        return page.get_text("text", clip=rect)

    def save_with_password(self, doc_id: str, output_path: Optional[str] = None, user_password: Optional[str] = None, owner_password: Optional[str] = None) -> Optional[str]:
        doc = self._acquire(doc_id)
        if not doc:
            return None
        save_path = output_path or doc.name
        try:
            import tempfile
            dir_name = os.path.dirname(os.path.abspath(save_path))
            fd, temp_path = tempfile.mkstemp(suffix='.pdf', dir=dir_name)
            os.close(fd)
            if user_password or owner_password:
                doc.save(temp_path, garbage=4, deflate=True, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=user_password or '', owner_pw=owner_password or user_password or '')
            else:
                doc.save(temp_path, garbage=4, deflate=True)
            if os.path.exists(save_path):
                os.replace(temp_path, save_path)
            else:
                os.rename(temp_path, save_path)
            self._dirty[doc_id] = False
            return save_path
        except Exception:
            return None

    def duplicate_page(self, doc_id: str, page_num: int) -> bool:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or page_num < 0 or page_num >= len(doc):
                return False
            doc.fullcopy_page(page_num, page_num + 1)
            info = self._infos.get(doc_id)
            if info:
                info.page_count = len(doc)
                info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
            self._dirty[doc_id] = True
            self._invalidate_render_cache(doc_id)
            return True

    def insert_blank_page(self, doc_id: str, index: int, width: float = 595, height: float = 842) -> bool:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return False
            idx = max(0, min(len(doc), index))
            doc.new_page(pno=idx, width=width, height=height)
            info = self._infos.get(doc_id)
            if info:
                info.page_count = len(doc)
                info.page_sizes = [PageSize(page_num=i, width=doc.load_page(i).rect.width, height=doc.load_page(i).rect.height) for i in range(len(doc))]
            self._dirty[doc_id] = True
            self._invalidate_render_cache(doc_id)
            return True

    def add_page_numbers(self, doc_id: str, prefix: str = "", start: int = 1, position: str = "bottom", fontsize: float = 10, color: str = "#000000") -> bool:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return False
            rgb = tuple(int(color.lstrip('#')[i:i+2], 16) / 255.0 for i in (0, 2, 4)) if color.startswith('#') else (0, 0, 0)
            n = len(doc)
            for i in range(n):
                page = doc.load_page(i)
                # Bates-style fixed-width number when a prefix is given, else "k / n"
                label = f"{prefix}{start + i:06d}" if prefix else f"{start + i} / {n}"
                rect = page.rect
                tw = fitz.get_text_length(label, fontsize=fontsize)
                x = rect.width / 2 - tw / 2
                y = rect.height - 18 if position == "bottom" else 24
                page.insert_text((x, y), label, fontsize=fontsize, color=rgb, overlay=True)
            self._dirty[doc_id] = True
            self._invalidate_render_cache(doc_id)
            return True

    def redact_matches(self, doc_id: str, query: str) -> int:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc or not query:
                return 0
            count = 0
            for i in range(len(doc)):
                page = doc.load_page(i)
                rects = page.search_for(query)
                for r in rects:
                    page.add_redact_annot(r, fill=(0, 0, 0))
                    count += 1
                if rects:
                    page.apply_redactions()
            if count:
                self._dirty[doc_id] = True
                self._invalidate_render_cache(doc_id)
            return count

    def generate_markup_summary(self, doc_id: str, annotations: List[Annotation]) -> Optional[dict]:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return None
            source_name = os.path.basename(doc.name)
        out = fitz.open()
        page = out.new_page()
        y = 60
        page.insert_text((50, y), f"Resumen de marcas — {source_name}", fontsize=15, color=(0, 0, 0))
        y += 28
        page.insert_text((50, y), f"Total: {len(annotations)} anotaciones", fontsize=10, color=(0.3, 0.3, 0.3))
        y += 22
        for idx, ann in enumerate(sorted(annotations, key=lambda a: (a.page, a.y or 0)), 1):
            if y > 780:
                page = out.new_page()
                y = 60
            text = (ann.text or "").replace("\n", " ").strip()
            line = f"{idx}.  Pág {ann.page + 1}  ·  {ann.type}" + (f"  ·  {text[:90]}" if text else "")
            page.insert_text((50, y), line, fontsize=9, color=(0.1, 0.1, 0.1))
            y += 15
        data = out.tobytes()
        out.close()
        return {
            "filename": source_name.replace('.pdf', '') + "_marcas.pdf",
            "data_base64": base64.b64encode(data).decode('utf-8'),
        }

    def compare_text(self, doc_id_a: str, doc_id_b: str) -> Optional[dict]:
        import difflib
        with self._lock:
            a = self._acquire(doc_id_a)
            b = self._acquire(doc_id_b)
            if not a or not b:
                return None
            text_a = [a.load_page(i).get_text() for i in range(len(a))]
            text_b = [b.load_page(i).get_text() for i in range(len(b))]
        diffs = []
        for i in range(max(len(text_a), len(text_b))):
            ta = text_a[i] if i < len(text_a) else ""
            tb = text_b[i] if i < len(text_b) else ""
            if ta == tb:
                continue
            sm = difflib.SequenceMatcher(None, ta.split(), tb.split())
            added, removed = [], []
            for tag, i1, i2, j1, j2 in sm.get_opcodes():
                if tag in ("replace", "delete"):
                    removed.extend(ta.split()[i1:i2])
                if tag in ("replace", "insert"):
                    added.extend(tb.split()[j1:j2])
            diffs.append({
                "page": i,
                "added": " ".join(added)[:400],
                "removed": " ".join(removed)[:400],
            })
        return {"pages_with_changes": len(diffs), "diffs": diffs}

    def close_document(self, doc_id: str) -> bool:
        with self._lock:
            existed = doc_id in self._infos
            doc = self._docs.pop(doc_id, None)
            if doc:
                doc.close()
            self._infos.pop(doc_id, None)
            self._dirty.pop(doc_id, None)
            self._passwords.pop(doc_id, None)
            self._lru.pop(doc_id, None)
            return existed

pdf_service = PdfService()
