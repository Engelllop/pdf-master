import fitz  # PyMuPDF
import base64
import logging
import math
import uuid
import json
import os
from typing import Dict, Optional, List
from collections import OrderedDict
from app.models.pdf import PdfInfo, PageRender, PdfOutlineItem, PageSize, Annotation
from app.core.config import settings
from app.services._pdf_base import PasswordRequiredError, DocumentNotFoundError

logger = logging.getLogger("pdfmaster")


class ExportMixin:
    def export_excel(self, doc_id: str, output_path: str) -> bool:
        with self._lock:
            doc = self._acquire(doc_id)
            if not doc:
                return False
            try:
                from openpyxl import Workbook
                wb = Workbook()
                first = True

                def next_sheet(title: str):
                    nonlocal first
                    ws = wb.active if first else wb.create_sheet()
                    ws.title = title[:31]
                    first = False
                    return ws

                for i in range(len(doc)):
                    page = doc.load_page(i)
                    tables = []
                    try:
                        tables = page.find_tables().tables
                    except Exception:
                        tables = []
                    if tables:
                        # Preserve real table structure, one sheet per table
                        for t_idx, tab in enumerate(tables):
                            ws = next_sheet(f"P{i+1}_T{t_idx+1}")
                            for row in tab.extract():
                                ws.append([("" if c is None else c) for c in row])
                    else:
                        # Fall back to raw text lines for pages without tables
                        ws = next_sheet(f"Página {i+1}")
                        for line in page.get_text().splitlines()[:1000]:
                            ws.append([line])
                self._guardar_atomico(output_path, False, lambda temp: wb.save(temp))
                return True
            except DocumentNotFoundError:
                raise
            except Exception:
                logger.exception("export_excel falló (doc %s)", doc_id)
                return False

    def export_pptx(self, doc_id: str, output_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        try:
            from pptx import Presentation
            from pptx.util import Inches
            from io import BytesIO
            prs = Presentation()
            blank_layout = prs.slide_layouts[6]
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(dpi=150)
                img_stream = BytesIO(pix.tobytes("png"))
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            self._guardar_atomico(output_path, False, lambda temp: prs.save(temp))
            return True
        except DocumentNotFoundError:
            raise
        except Exception:
            logger.exception("export_pptx falló (doc %s)", doc_id)
            return False

    def export_txt(self, doc_id: str, output_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        try:
            def escribir(temp: str) -> None:
                with open(temp, 'w', encoding='utf-8') as f:
                    for i in range(len(doc)):
                        f.write(doc.load_page(i).get_text())
                        if i < len(doc) - 1:
                            f.write('\n\f\n')  # form feed entre páginas

            # `open(output_path, 'w')` truncaba el archivo del usuario ANTES de escribir
            # una sola letra: si fallaba en la página 1, donde había un documento
            # quedaba un archivo vacío.
            self._guardar_atomico(output_path, False, escribir)
            return True
        except DocumentNotFoundError:
            raise
        except Exception:
            logger.exception("export_txt falló (doc %s)", doc_id)
            return False

    def export_html(self, doc_id: str, output_path: str) -> bool:
        doc = self._acquire(doc_id)
        if not doc:
            return False
        try:
            parts = ['<!DOCTYPE html><html><head><meta charset="utf-8"></head><body>']
            for i in range(len(doc)):
                parts.append(doc.load_page(i).get_text("html"))
            parts.append('</body></html>')
            def escribir(temp: str) -> None:
                with open(temp, 'w', encoding='utf-8') as f:
                    f.write('\n'.join(parts))

            self._guardar_atomico(output_path, False, escribir)
            return True
        except DocumentNotFoundError:
            raise
        except Exception:
            logger.exception("export_html falló (doc %s)", doc_id)
            return False

    def export_word(self, doc_id: str, output_path: Optional[str] = None) -> Optional[dict]:
        """Con output_path escribe donde el usuario eligió (como Excel y PowerPoint);
        sin él devuelve base64, como antes."""
        doc = self._acquire(doc_id)
        if not doc:
            return None
        try:
            from docx import Document
            from docx.shared import Pt
            import io
            import base64
            document = Document()
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    for line in text.split('\n'):
                        if line.strip():
                            p = document.add_paragraph(line.strip())
                            p.style.font.size = Pt(11)
                if page_num < len(doc) - 1:
                    document.add_page_break()
            filename = os.path.basename(self._doc_path(doc_id)).replace('.pdf', '.docx')
            if output_path:
                self._guardar_atomico(output_path, False, lambda temp: document.save(temp))
                return {"filename": filename, "output_path": output_path}
            buffer = io.BytesIO()
            document.save(buffer)
            buffer.seek(0)
            data = buffer.read()
            return {
                "filename": filename,
                "data_base64": base64.b64encode(data).decode('utf-8'),
            }
        except DocumentNotFoundError:
            raise
        except Exception:
            logger.exception("export_word falló (doc %s)", doc_id)
            return None

    def export_measurements(self, rows: List[dict], output_path: str, title: str = "") -> bool:
        """Tabla de mediciones/conteos a CSV o XLSX. No toca fitz (sin lock)."""
        # `Escala` por fila: un juego de planos mezcla escalas y una sola en el título
        # no dice con cuál se tomó cada cota.
        headers = ["Página", "Tipo", "Etiqueta", "Valor", "Unidad", "Escala"]
        try:
            if output_path.lower().endswith(".csv"):
                import csv

                def escribir(temp: str) -> None:
                    # utf-8-sig para que Excel detecte acentos al abrir el CSV
                    with open(temp, "w", newline="", encoding="utf-8-sig") as f:
                        w = csv.writer(f, delimiter=";")
                        if title:
                            w.writerow([title])
                            w.writerow([])
                        w.writerow(headers)
                        for r in rows:
                            w.writerow([r.get("page"), r.get("tipo"), r.get("etiqueta"),
                                        r.get("valor"), r.get("unidad"), r.get("escala", "")])

                self._guardar_atomico(output_path, False, escribir)
            else:
                from openpyxl import Workbook
                wb = Workbook()
                ws = wb.active
                ws.title = "Mediciones"
                row_idx = 1
                if title:
                    ws.cell(row=row_idx, column=1, value=title)
                    row_idx += 2
                for col, h in enumerate(headers, start=1):
                    ws.cell(row=row_idx, column=col, value=h)
                for r in rows:
                    row_idx += 1
                    ws.cell(row=row_idx, column=1, value=r.get("page"))
                    ws.cell(row=row_idx, column=2, value=r.get("tipo"))
                    ws.cell(row=row_idx, column=3, value=r.get("etiqueta"))
                    ws.cell(row=row_idx, column=4, value=r.get("valor"))
                    ws.cell(row=row_idx, column=5, value=r.get("unidad"))
                    ws.cell(row=row_idx, column=6, value=r.get("escala", ""))
                self._guardar_atomico(output_path, False, lambda temp: wb.save(temp))
            return True
        except Exception:
            logger.exception("export_measurements falló (%s)", output_path)
            return False
