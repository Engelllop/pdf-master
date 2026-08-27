"""Cada página se pegaba en la diapositiva a 150 dpi fijos y estirada al tamaño exacto
de la diapositiva: un plano de 36×24 in salía a 5400×3600 px (trece veces más píxeles
de los que la diapositiva muestra) y encima achatado."""
import io

import fitz
import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Emu


@pytest.fixture
def plano_grande(tmp_path):
    """Dos láminas apaisadas de 36×24 in (2592×1728 pt)."""
    ruta = tmp_path / "planos.pdf"
    doc = fitz.open()
    for i in range(2):
        pg = doc.new_page(width=2592, height=1728)
        pg.insert_text((100, 100), f"LAMINA A-10{i + 1}", fontsize=48)
    doc.save(str(ruta))
    doc.close()
    return str(ruta)


def _exportar(client, ruta, salida):
    doc_id = client.post("/pdf/open", json={"file_path": ruta}).json()["doc_id"]
    r = client.post(f"/pdf/export-pptx/{doc_id}?output_path={salida}")
    client.post(f"/pdf/close/{doc_id}")
    assert r.status_code == 200, r.text
    return Presentation(salida)


def _imagenes(prs):
    """(shape, ancho_px, alto_px) de cada imagen incrustada."""
    salida = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                with Image.open(io.BytesIO(shape.image.blob)) as im:
                    salida.append((shape, im.width, im.height))
    return salida


class TestExportPptx:
    def test_el_bitmap_no_pasa_del_tope(self, client, plano_grande, tmp_path):
        salida = str(tmp_path / "a.pptx")
        prs = _exportar(client, plano_grande, salida)
        imagenes = _imagenes(prs)
        assert len(imagenes) == 2
        for _, w, h in imagenes:
            assert max(w, h) <= 2000, f"{w}x{h} px: sin tope, un plano sale a 5400 px"

    def test_no_se_deforma_la_lamina(self, client, plano_grande, tmp_path):
        salida = str(tmp_path / "b.pptx")
        prs = _exportar(client, plano_grande, salida)
        shape, w, h = _imagenes(prs)[0]
        # La proporción de la forma en la diapositiva es la del bitmap (±1 %).
        assert shape.width / shape.height == pytest.approx(w / h, rel=0.01)

    def test_queda_centrada_y_dentro(self, client, plano_grande, tmp_path):
        salida = str(tmp_path / "c.pptx")
        prs = _exportar(client, plano_grande, salida)
        shape, _, _ = _imagenes(prs)[0]
        assert shape.left >= 0 and shape.top >= 0
        assert shape.left + shape.width <= prs.slide_width + Emu(1)
        assert shape.top + shape.height <= prs.slide_height + Emu(1)
        # Centrada: los márgenes de un lado y del otro son iguales.
        assert shape.left == pytest.approx((prs.slide_width - shape.width) / 2, abs=2)
        assert shape.top == pytest.approx((prs.slide_height - shape.height) / 2, abs=2)

    def test_una_pagina_chica_no_se_estira_a_lo_bestia(self, client, tmp_path):
        # Una A4 vertical: la proporción también se respeta (antes se estiraba al 4:3).
        ruta = tmp_path / "a4.pdf"
        doc = fitz.open()
        doc.new_page(width=595, height=842)
        doc.save(str(ruta))
        doc.close()
        prs = _exportar(client, str(ruta), str(tmp_path / "d.pptx"))
        shape, w, h = _imagenes(prs)[0]
        assert h > w  # sigue siendo vertical
        assert shape.height > shape.width
